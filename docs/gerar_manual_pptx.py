# -*- coding: utf-8 -*-
"""Gera a apresentação (manual em slides) do sistema Escala em formato .pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PRETO = RGBColor(0x1A, 0x1A, 0x1A)
PRETO_FORTE = RGBColor(0x00, 0x00, 0x00)
CINZA = RGBColor(0x4B, 0x55, 0x63)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF6, 0xF8)

LARGURA = Inches(13.333)
ALTURA = Inches(7.5)


def novo_slide_branco(prs):
    layout_branco = prs.slide_layouts[6]
    return prs.slides.add_slide(layout_branco)


def add_fundo(slide, cor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, LARGURA, ALTURA)
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor
    shape.line.fill.background()
    shape.shadow.inherit = False
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    return shape


def add_faixa_topo(slide, altura=Inches(1.0), cor=PRETO):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, LARGURA, altura)
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_titulo_slide(slide, texto, top=Inches(0.15), tamanho=32, cor=BRANCO, largura=None, esquerda=Inches(0.6)):
    largura = largura or (LARGURA - Inches(1.2))
    box = slide.shapes.add_textbox(esquerda, top, largura, Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.bold = True
    run.font.color.rgb = cor
    return box


def add_rodape(slide, numero, total):
    box = slide.shapes.add_textbox(LARGURA - Inches(1.5), ALTURA - Inches(0.5), Inches(1.2), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f'{numero} / {total}'
    run.font.size = Pt(11)
    run.font.color.rgb = CINZA
    box2 = slide.shapes.add_textbox(Inches(0.6), ALTURA - Inches(0.5), Inches(4), Inches(0.35))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = 'Sistema Escala — Manual do Usuário'
    run2.font.size = Pt(11)
    run2.font.color.rgb = CINZA


def add_bullets(slide, itens, top=Inches(1.3), esquerda=Inches(0.6), largura=None, altura=None, tamanho=18, cor=PRETO, espaco_depois=Pt(10)):
    largura = largura or (LARGURA - Inches(1.2))
    altura = altura or (ALTURA - top - Inches(0.6))
    box = slide.shapes.add_textbox(esquerda, top, largura, altura)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = espaco_depois
        if isinstance(item, tuple):
            nivel, texto_item, negrito_prefixo = item
            p.level = nivel
            if negrito_prefixo:
                r1 = p.add_run()
                r1.text = negrito_prefixo
                r1.font.bold = True
                r1.font.size = Pt(tamanho - (2 * nivel))
                r1.font.color.rgb = cor
                r2 = p.add_run()
                r2.text = texto_item
                r2.font.size = Pt(tamanho - (2 * nivel))
                r2.font.color.rgb = cor
            else:
                r = p.add_run()
                r.text = '• ' + texto_item if nivel == 0 else '– ' + texto_item
                r.font.size = Pt(tamanho - (2 * nivel))
                r.font.color.rgb = cor
        else:
            r = p.add_run()
            r.text = '• ' + item
            r.font.size = Pt(tamanho)
            r.font.color.rgb = cor
    return box


def add_tabela(slide, cabecalhos, linhas, top=Inches(1.5), esquerda=Inches(0.6), largura=None, altura=Inches(4.5)):
    largura = largura or (LARGURA - Inches(1.2))
    n_linhas = len(linhas) + 1
    n_colunas = len(cabecalhos)
    grafico = slide.shapes.add_table(n_linhas, n_colunas, esquerda, top, largura, altura)
    tabela = grafico.table
    for i, titulo in enumerate(cabecalhos):
        celula = tabela.cell(0, i)
        celula.text = titulo
        celula.fill.solid()
        celula.fill.fore_color.rgb = PRETO
        for p in celula.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = BRANCO
    for r, linha in enumerate(linhas, start=1):
        for c, valor in enumerate(linha):
            celula = tabela.cell(r, c)
            celula.text = str(valor)
            for p in celula.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = PRETO
    return grafico


def slide_secao(prs, titulo, subtitulo=None):
    slide = novo_slide_branco(prs)
    add_fundo(slide, PRETO)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), LARGURA - Inches(1.6), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = titulo
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = BRANCO
    if subtitulo:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitulo
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    return slide


def slide_conteudo(prs, titulo):
    slide = novo_slide_branco(prs)
    add_fundo(slide, BRANCO)
    add_faixa_topo(slide)
    add_titulo_slide(slide, titulo)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA

    # 1. Capa
    slide = novo_slide_branco(prs)
    add_fundo(slide, PRETO)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), LARGURA - Inches(1.6), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Manual do Usuário'
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = BRANCO

    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), LARGURA - Inches(1.6), Inches(0.8))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = 'Sistema Escala'
    run2.font.size = Pt(28)
    run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    box3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), LARGURA - Inches(1.6), Inches(0.6))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = 'Gestão de voluntários, ministérios, eventos e escalas de serviço'
    run3.font.size = Pt(16)
    run3.italic = True
    run3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # 2. Agenda
    slide = slide_conteudo(prs, 'Agenda')
    add_bullets(slide, [
        'O que é o sistema Escala',
        'Perfis de usuário: Voluntário e Administrador',
        'Login e cadastro',
        'Área do Voluntário',
        'Área do Administrador',
        'Notificações via WhatsApp',
        'Perguntas frequentes',
    ])

    # 3. O que é o sistema
    slide = slide_conteudo(prs, 'O que é o sistema Escala')
    add_bullets(slide, [
        'Organiza a escala de voluntários de uma igreja/organização.',
        'Permite cadastrar ministérios (equipes de serviço).',
        'Permite cadastrar eventos (cultos, reuniões) com necessidades de voluntários por ministério.',
        'Gera escalas automaticamente, respeitando indisponibilidades cadastradas.',
        'Notifica voluntários via WhatsApp e recebe a confirmação/recusa deles.',
    ])

    # 4. Seção: Perfis
    slide_secao(prs, 'Perfis de Usuário')

    # 5. Perfis - tabela
    slide = slide_conteudo(prs, 'Voluntário x Administrador')
    add_tabela(
        slide,
        ['Aspecto', 'Voluntário', 'Administrador'],
        [
            ['Tela inicial', 'Minha Escala', 'Escala (admin)'],
            ['Menu', 'Minha Escala, Indisponibilidade', 'Escala, Eventos, Voluntários, Ministérios'],
            ['Autocadastro', 'Sim (tela de Registro)', 'Não; criado por um administrador'],
            ['Gera escalas', 'Não', 'Sim, manual ou automática'],
            ['Responde escala', 'Confirma/recusa a própria', 'Pode confirmar/recusar de qualquer um'],
        ],
        top=Inches(1.4),
        altura=Inches(4.6),
    )

    # 6. Seção: Login e cadastro
    slide_secao(prs, 'Login e Cadastro')

    # 7. Login
    slide = slide_conteudo(prs, 'Tela de Login')
    add_bullets(slide, [
        'Disponível em /login.',
        (0, 'Campos: E-mail e Senha.', None),
        (0, 'Clique em "Entrar" para acessar o sistema.', None),
        (0, 'Administradores são levados à tela de Escala; voluntários, à tela Minha Escala.', None),
        (0, 'Link "Cadastre-se como voluntário" leva à tela de Registro.', None),
    ])

    # 8. Registro
    slide = slide_conteudo(prs, 'Cadastro de Novo Voluntário')
    add_bullets(slide, [
        'Disponível em /registrar, para autoatendimento de novos voluntários.',
        (0, 'Campos: Nome completo, E-mail, Telefone/WhatsApp, Senha e Ministérios de interesse.', None),
        (0, 'Após cadastrar, o login é feito automaticamente.', None),
        (0, 'Todo cadastro feito aqui cria um usuário do tipo Voluntário.', None),
    ])

    # 9. Seção: Área do voluntário
    slide_secao(prs, 'Área do Voluntário')

    # 10. Minha Escala
    slide = slide_conteudo(prs, 'Minha Escala')
    add_bullets(slide, [
        'Lista as escalas (atribuições) futuras do voluntário logado.',
        (0, 'Colunas: Evento, Data, Ministério, Status e Ações.', None),
        (0, 'Status: Pendente, Confirmado ou Recusado.', None),
        (0, 'Para escalas pendentes: botões "Confirmar" e "Não posso".', None),
    ])

    # 11. Indisponibilidade
    slide = slide_conteudo(prs, 'Indisponibilidade')
    add_bullets(slide, [
        'Permite informar datas em que o voluntário não poderá servir.',
        (0, 'Campos: Data (obrigatório) e Motivo (opcional).', None),
        (0, 'Clique em "Adicionar" para cadastrar.', None),
        (0, 'A lista abaixo permite "Remover" qualquer indisponibilidade cadastrada.', None),
        (0, 'Datas indisponíveis não são usadas na geração automática de escala.', None),
    ])

    # 12. Seção: Área do administrador
    slide_secao(prs, 'Área do Administrador')

    # 13. Ministérios
    slide = slide_conteudo(prs, 'Ministérios')
    add_bullets(slide, [
        'Cadastro das equipes de serviço (ex.: Louvor, Recepção, Mídia).',
        (0, 'Formulário simples: nome do ministério + botão "Adicionar".', None),
        (0, 'Botão "Remover" exclui o ministério (após confirmação).', None),
    ])

    # 14. Voluntários
    slide = slide_conteudo(prs, 'Voluntários')
    add_bullets(slide, [
        'Cadastro de novos voluntários pelo administrador.',
        (0, 'Campos: Nome, E-mail, Telefone, Senha provisória e Ministérios.', None),
        'Gestão dos voluntários existentes:',
        (1, '"Editar ministérios" — altera os ministérios do voluntário.', None),
        (1, '"Ativar/Desativar" — controla se ele participa da geração de escala.', None),
        (1, '"Remover" — exclui o voluntário.', None),
    ])

    # 15. Eventos
    slide = slide_conteudo(prs, 'Eventos')
    add_bullets(slide, [
        'Cadastro de cultos/reuniões com necessidades de voluntários.',
        (0, 'Campos: Nome do evento e Data/Hora.', None),
        (0, 'Necessidades: selecionar Ministério + quantidade de vagas (pode adicionar várias).', None),
        (0, 'Botão "Criar evento" salva o evento com suas necessidades.', None),
        (0, 'Botão "Remover" exclui o evento e sua escala associada.', None),
    ])

    # 16. Escala - geração
    slide = slide_conteudo(prs, 'Escala — Geração Automática')
    add_bullets(slide, [
        'Tela principal do administrador, em /admin/escala.',
        (0, '"Gerar escala automática (próximos 30 dias)" — distribui voluntários em todos os eventos do período.', None),
        (0, '"Gerar/completar escala deste evento" — preenche as vagas em falta de um evento específico.', None),
        (0, 'A geração respeita as indisponibilidades cadastradas pelos voluntários.', None),
        (0, 'Cada voluntário escalado recebe uma notificação automática via WhatsApp.', None),
    ])

    # 17. Escala - gestão
    slide = slide_conteudo(prs, 'Escala — Gestão das Atribuições')
    add_bullets(slide, [
        'Para cada escala gerada, é possível:',
        (1, 'Reatribuir o voluntário (lista suspensa com voluntários elegíveis).', None),
        (1, '"Confirmar" / "Recusar" manualmente.', None),
        (1, '"Reenviar notificação" via WhatsApp.', None),
        (1, '"Remover" a atribuição.', None),
        (0, 'O selo "Notificado" mostra se e quando a mensagem foi enviada.', None),
    ])

    # 18. Seção: WhatsApp
    slide_secao(prs, 'Notificações via WhatsApp')

    # 19. WhatsApp detalhe
    slide = slide_conteudo(prs, 'Como funcionam as notificações')
    add_bullets(slide, [
        'Ao ser escalado, o voluntário recebe uma mensagem com ministério, evento e data.',
        (0, 'Basta responder "SIM" para confirmar ou "NAO" para recusar, direto pelo WhatsApp.', None),
        (0, 'Os envios em lote são feitos de forma sequencial, com um intervalo entre cada mensagem.', None),
        (0, 'Isso reduz o risco de bloqueio do número ao notificar muitos voluntários de uma vez.', None),
        (0, 'O status do envio pode ser consultado na coluna "Notificado" da tela de Escala.', None),
    ])

    # 20. FAQ
    slide = slide_conteudo(prs, 'Perguntas Frequentes')
    add_bullets(slide, [
        (0, 'Não recebi a notificação: verifique seu telefone cadastrado e peça reenvio ao administrador.', None),
        (0, 'Não posso em uma data: cadastre-a em Indisponibilidade antes da geração da escala.', None),
        (0, 'Confirmei mas não poderei mais ir: avise o administrador para reatribuir a vaga.', None),
        (0, 'Esqueci minha senha: solicite a um administrador a atualização do seu acesso.', None),
        (0, 'Posso servir em mais de um ministério: sim, é possível selecionar vários.', None),
    ], tamanho=16)

    # 21. Encerramento
    slide = novo_slide_branco(prs)
    add_fundo(slide, PRETO)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), LARGURA - Inches(1.6), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'Obrigado!'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = BRANCO
    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), LARGURA - Inches(1.6), Inches(0.8))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = 'Sistema Escala — Manual do Usuário'
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    caminho = 'docs/Manual_do_Usuario_Escala.pptx'
    prs.save(caminho)
    print('Gerado:', caminho)


if __name__ == '__main__':
    build()
